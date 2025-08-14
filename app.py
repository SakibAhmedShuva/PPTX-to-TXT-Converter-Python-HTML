from flask import Flask, request, jsonify, render_template, send_from_directory
from flask_cors import CORS
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
import os
import subprocess
from werkzeug.utils import secure_filename

app = Flask(__name__)
CORS(app)  # Enable CORS for all routes

# --- Configuration ---
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pptx', 'ppt'}
MAX_CONTENT_LENGTH = 16 * 1024 * 1024  # 16MB max file size

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

# Create upload directory if it doesn't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)


# --- Helper Functions ---
def allowed_file(filename):
    """Check if the file extension is allowed."""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pptx(file_path):
    """
    Extract text from PowerPoint, recursively handling grouped shapes and using
    absolute positioning for correct sorting.
    """
    try:
        presentation = Presentation(file_path)
        text_content = []

        def get_shapes_with_abs_position(shapes, parent_top=0, parent_left=0):
            shape_info_list = []
            for shape in shapes:
                shape_top = shape.top.emu if hasattr(shape, 'top') and shape.top is not None else 0
                shape_left = shape.left.emu if hasattr(shape, 'left') and shape.left is not None else 0
                abs_top = parent_top + shape_top
                abs_left = parent_left + shape_left
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    shape_info_list.extend(get_shapes_with_abs_position(shape.shapes, abs_top, abs_left))
                else:
                    shape_info_list.append({'shape': shape, 'top': abs_top, 'left': abs_left})
            return shape_info_list

        for i, slide in enumerate(presentation.slides, 1):
            slide_text = f"=== SLIDE {i} ===\n"
            all_shape_infos = get_shapes_with_abs_position(slide.shapes)
            added_placeholder_idxs = {
                info['shape'].placeholder_format.idx
                for info in all_shape_infos
                if info['shape'].is_placeholder
            }
            for shape in slide.slide_layout.placeholders:
                if shape.placeholder_format.idx not in added_placeholder_idxs:
                    top = shape.top.emu if hasattr(shape, 'top') and shape.top is not None else 0
                    left = shape.left.emu if hasattr(shape, 'left') and shape.left is not None else 0
                    all_shape_infos.append({'shape': shape, 'top': top, 'left': left})
                    added_placeholder_idxs.add(shape.placeholder_format.idx)
            for shape in slide.slide_layout.slide_master.placeholders:
                 if shape.placeholder_format.idx not in added_placeholder_idxs:
                    top = shape.top.emu if hasattr(shape, 'top') and shape.top is not None else 0
                    left = shape.left.emu if hasattr(shape, 'left') and shape.left is not None else 0
                    all_shape_infos.append({'shape': shape, 'top': top, 'left': left})
            text_shape_infos = [info for info in all_shape_infos if hasattr(info['shape'], "text")]
            sorted_shape_infos = sorted(text_shape_infos, key=lambda info: (info['top'], info['left']))
            for info in sorted_shape_infos:
                shape = info['shape']
                text = ""
                if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                    text = str(i)
                elif hasattr(shape, "text"):
                    text = shape.text.strip()
                if text:
                    slide_text += text + "\n"
            text_content.append(slide_text)
        return "\n".join(text_content)
    except Exception as e:
        raise Exception(f"Error processing PowerPoint file: {str(e)}")


# --- Routes ---
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/convert-to-text', methods=['POST'])
def convert_pptx_to_text():
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    file = request.files['file']
    if file.filename == '' or not allowed_file(file.filename):
        return jsonify({'error': 'Invalid file or file type'}), 400
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    try:
        file.save(filepath)
        extracted_text = extract_text_from_pptx(filepath)
        return jsonify({'success': True, 'text': extracted_text, 'filename': filename}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        if os.path.exists(filepath):
            os.remove(filepath)

@app.route('/api/convert-to-pdf', methods=['POST'])
def convert_pptx_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    file = request.files['file']
    if file.filename == '' or not allowed_file(file.filename):
        return jsonify({'error': 'Invalid file or file type'}), 400
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    pdf_filename = filename.rsplit('.', 1)[0] + '.pdf'
    pdf_filepath = os.path.join(app.config['UPLOAD_FOLDER'], pdf_filename)
    try:
        file.save(filepath)
        subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', app.config['UPLOAD_FOLDER'], filepath],
            check=True, timeout=60
        )
        if not os.path.exists(pdf_filepath):
             raise FileNotFoundError("PDF conversion failed, output file not found.")
        # MODIFIED: Removed as_attachment=True to allow inline viewing
        return send_from_directory(app.config['UPLOAD_FOLDER'], pdf_filename)
    except subprocess.TimeoutExpired:
        return jsonify({'error': 'PDF conversion timed out after 60 seconds.'}), 500
    except (subprocess.CalledProcessError, FileNotFoundError, Exception) as e:
        return jsonify({'error': f'PDF conversion failed: {str(e)}'}), 500
    finally:
        if os.path.exists(filepath):
            os.remove(filepath)
        if os.path.exists(pdf_filepath):
            os.remove(pdf_filepath)

# --- Error Handlers ---
@app.errorhandler(413)
def too_large(e):
    return jsonify({'error': 'File too large. Maximum size is 16MB'}), 413

# --- Main Execution ---
if __name__ == '__main__':
    print("🚀 Starting PPTX Converter API...")
    print("📝 Interface available at: http://localhost:7860")
    app.run(host='0.0.0.0', port=7860, debug=True)